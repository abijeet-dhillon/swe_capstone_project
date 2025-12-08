"""
Document Parser Module
Extracts text content from various document formats.

Supported formats:
- Documents: .docx, .pdf, .rtf, .txt, .md
- Presentations: .pptx
- Spreadsheets: .csv, .xlsx, .xls
- Code: .py, .js, .java, .cpp, .c, .ts, .jsx, .tsx, .go, .rs, .rb, .php, .swift, .kt
- Markup: .html, .xml, .json, .yaml, .yml
"""
from pathlib import Path
from typing import Dict, Any, List, Optional
import mimetypes

# Document parsing
from docx import Document
from pptx import Presentation

# PDF parsing
try:
    import pypdf
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

# Excel parsing
try:
    import openpyxl
    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False

# HTML parsing
try:
    from bs4 import BeautifulSoup
    HAS_HTML = True
except ImportError:
    HAS_HTML = False

# RTF parsing
try:
    from striprtf.striprtf import rtf_to_text
    HAS_RTF = True
except ImportError:
    HAS_RTF = False

import csv
import json
import yaml


class DocumentParser:
    """Parser for extracting text from various document formats."""
    
    # Define supported extensions
    DOCUMENT_FORMATS = {'.docx', '.pdf', '.rtf', '.txt', '.md', '.rst'}
    PRESENTATION_FORMATS = {'.pptx'}
    SPREADSHEET_FORMATS = {'.csv', '.xlsx', '.xls'}
    CODE_FORMATS = {'.py', '.js', '.java', '.cpp', '.c', '.h', '.hpp', '.ts', '.jsx', 
                    '.tsx', '.go', '.rs', '.rb', '.php', '.swift', '.kt', '.cs', '.r', '.m'}
    MARKUP_FORMATS = {'.html', '.htm', '.xml', '.json', '.yaml', '.yml'}
    
    @staticmethod
    def get_supported_formats() -> List[str]:
        """Return list of all supported file formats."""
        return sorted(
            DocumentParser.DOCUMENT_FORMATS |
            DocumentParser.PRESENTATION_FORMATS |
            DocumentParser.SPREADSHEET_FORMATS |
            DocumentParser.CODE_FORMATS |
            DocumentParser.MARKUP_FORMATS
        )
    
    @staticmethod
    def is_supported(file_path: str) -> bool:
        """Check if file format is supported."""
        suffix = Path(file_path).suffix.lower()
        return suffix in DocumentParser.get_supported_formats()
    
    # ==================== DOCUMENT PARSERS ====================
    
    @staticmethod
    def parse_docx(file_path: str) -> Dict[str, Any]:
        """Extract text content from a Word document (.docx)."""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        try:
            doc = Document(file_path)
            
            # Extract all text from paragraphs
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            full_text = "\n\n".join(paragraphs)
            
            # Extract text from tables
            table_text = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        table_text.append(row_text)
            
            if table_text:
                full_text += "\n\nTables:\n" + "\n".join(table_text)
            
            return {
                "file_name": file_path.name,
                "file_type": "docx",
                "text": full_text,
                "paragraph_count": len(paragraphs),
                "table_count": len(doc.tables),
                "word_count": len(full_text.split()),
                "metadata": {}
            }
            
        except Exception as e:
            raise RuntimeError(f"Failed to parse DOCX file: {str(e)}")
    
    @staticmethod
    def parse_pdf(file_path: str) -> Dict[str, Any]:
        """Extract text content from a PDF document."""
        if not HAS_PDF:
            raise RuntimeError("PDF support not installed. Run: pip install pypdf")
        
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        try:
            with open(file_path, 'rb') as f:
                pdf = pypdf.PdfReader(f)
                
                # Extract text from all pages
                pages_text = []
                for page_num, page in enumerate(pdf.pages, start=1):
                    text = page.extract_text()
                    if text.strip():
                        pages_text.append(f"--- Page {page_num} ---\n{text}")
                
                full_text = "\n\n".join(pages_text)
                
                # Extract metadata
                metadata = {}
                if pdf.metadata:
                    metadata = {
                        "title": pdf.metadata.get('/Title', ''),
                        "author": pdf.metadata.get('/Author', ''),
                        "subject": pdf.metadata.get('/Subject', ''),
                        "creator": pdf.metadata.get('/Creator', '')
                    }
                
                return {
                    "file_name": file_path.name,
                    "file_type": "pdf",
                    "text": full_text,
                    "page_count": len(pdf.pages),
                    "word_count": len(full_text.split()),
                    "metadata": metadata
                }
                
        except Exception as e:
            raise RuntimeError(f"Failed to parse PDF file: {str(e)}")
    
    @staticmethod
    def parse_txt(file_path: str) -> Dict[str, Any]:
        """Extract text content from plain text files (.txt, .md, .rst, etc.)."""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252']
            text = None
            used_encoding = 'utf-8'
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text = f.read()
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if text is None:
                raise RuntimeError("Could not decode file with any supported encoding")
            
            lines = text.split('\n')
            non_empty_lines = [line for line in lines if line.strip()]
            
            return {
                "file_name": file_path.name,
                "file_type": file_path.suffix.lstrip('.') or "txt",
                "text": text,
                "line_count": len(lines),
                "non_empty_lines": len(non_empty_lines),
                "word_count": len(text.split()),
                "encoding": used_encoding,
                "metadata": {}
            }
            
        except Exception as e:
            raise RuntimeError(f"Failed to parse text file: {str(e)}")
    
    @staticmethod
    def parse_rtf(file_path: str) -> Dict[str, Any]:
        """Extract text content from RTF files."""
        if not HAS_RTF:
            raise RuntimeError("RTF support not installed. Run: pip install striprtf")
        
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            
            text = rtf_to_text(rtf_content)
            
            return {
                "file_name": file_path.name,
                "file_type": "rtf",
                "text": text,
                "word_count": len(text.split()),
                "metadata": {}
            }
            
        except Exception as e:
            raise RuntimeError(f"Failed to parse RTF file: {str(e)}")
    
    # ==================== PRESENTATION PARSERS ====================
    
    @staticmethod
    def parse_pptx(file_path: str) -> Dict[str, Any]:
        """Extract text content from a PowerPoint presentation (.pptx)."""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        try:
            prs = Presentation(file_path)
            
            slides_text = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_content = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                if slide_content:
                    slide_text = f"Slide {slide_num}:\n" + "\n".join(slide_content)
                    slides_text.append(slide_text)
            
            full_text = "\n\n".join(slides_text)
            
            return {
                "file_name": file_path.name,
                "file_type": "pptx",
                "text": full_text,
                "slide_count": len(prs.slides),
                "word_count": len(full_text.split()),
                "metadata": {}
            }
            
        except Exception as e:
            raise RuntimeError(f"Failed to parse PPTX file: {str(e)}")
    
    # ==================== SPREADSHEET PARSERS ====================
    
    @staticmethod
    def parse_csv(file_path: str) -> Dict[str, Any]:
        """Extract text content from CSV files."""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        try:
            rows = []
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                for row in reader:
                    row_text = " | ".join(str(cell).strip() for cell in row if str(cell).strip())
                    if row_text:
                        rows.append(row_text)
            
            full_text = "\n".join(rows)
            
            return {
                "file_name": file_path.name,
                "file_type": "csv",
                "text": full_text,
                "row_count": len(rows),
                "word_count": len(full_text.split()),
                "metadata": {}
            }
            
        except Exception as e:
            raise RuntimeError(f"Failed to parse CSV file: {str(e)}")
    
    @staticmethod
    def parse_excel(file_path: str) -> Dict[str, Any]:
        """Extract text content from Excel files (.xlsx, .xls)."""
        if not HAS_EXCEL:
            raise RuntimeError("Excel support not installed. Run: pip install openpyxl")
        
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            
            sheets_text = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_rows = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell).strip() for cell in row if cell is not None and str(cell).strip())
                    if row_text:
                        sheet_rows.append(row_text)
                
                if sheet_rows:
                    sheet_text = f"Sheet: {sheet_name}\n" + "\n".join(sheet_rows)
                    sheets_text.append(sheet_text)
            
            full_text = "\n\n".join(sheets_text)
            
            return {
                "file_name": file_path.name,
                "file_type": "xlsx" if file_path.suffix == '.xlsx' else "xls",
                "text": full_text,
                "sheet_count": len(wb.sheetnames),
                "word_count": len(full_text.split()),
                "metadata": {}
            }
            
        except Exception as e:
            raise RuntimeError(f"Failed to parse Excel file: {str(e)}")
    
    # ==================== CODE PARSERS ====================
    
    @staticmethod
    def parse_code(file_path: str) -> Dict[str, Any]:
        """Extract content from code files."""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252']
            code = None
            used_encoding = 'utf-8'
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        code = f.read()
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if code is None:
                raise RuntimeError("Could not decode file with any supported encoding")
            
            lines = code.split('\n')
            non_empty_lines = [line for line in lines if line.strip()]
            
            # Count comments (simple heuristic)
            comment_markers = {
                '.py': '#',
                '.js': '//',
                '.ts': '//',
                '.java': '//',
                '.cpp': '//',
                '.c': '//',
                '.go': '//',
                '.rs': '//',
                '.rb': '#',
                '.php': '//',
                '.swift': '//',
                '.kt': '//'
            }
            
            comment_char = comment_markers.get(file_path.suffix.lower(), '#')
            comment_lines = [line for line in non_empty_lines if line.strip().startswith(comment_char)]
            
            return {
                "file_name": file_path.name,
                "file_type": file_path.suffix.lstrip('.'),
                "text": code,
                "line_count": len(lines),
                "non_empty_lines": len(non_empty_lines),
                "comment_lines": len(comment_lines),
                "word_count": len(code.split()),
                "encoding": used_encoding,
                "language": file_path.suffix.lstrip('.'),
                "metadata": {}
            }
            
        except Exception as e:
            raise RuntimeError(f"Failed to parse code file: {str(e)}")
    
    # ==================== MARKUP PARSERS ====================
    
    @staticmethod
    def parse_html(file_path: str) -> Dict[str, Any]:
        """Extract text content from HTML files."""
        if not HAS_HTML:
            raise RuntimeError("HTML support not installed. Run: pip install beautifulsoup4 lxml")
        
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'lxml')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            # Extract title
            title = soup.title.string if soup.title else ""
            
            return {
                "file_name": file_path.name,
                "file_type": "html",
                "text": text,
                "word_count": len(text.split()),
                "title": title,
                "metadata": {"title": title}
            }
            
        except Exception as e:
            raise RuntimeError(f"Failed to parse HTML file: {str(e)}")
    
    @staticmethod
    def parse_json(file_path: str) -> Dict[str, Any]:
        """Extract content from JSON files."""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Convert to pretty-printed text
            text = json.dumps(data, indent=2)
            
            return {
                "file_name": file_path.name,
                "file_type": "json",
                "text": text,
                "word_count": len(text.split()),
                "metadata": {}
            }
            
        except Exception as e:
            raise RuntimeError(f"Failed to parse JSON file: {str(e)}")
    
    @staticmethod
    def parse_yaml(file_path: str) -> Dict[str, Any]:
        """Extract content from YAML files."""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                yaml_content = f.read()
                data = yaml.safe_load(yaml_content)
            
            # Convert back to text (use original + parsed for validation)
            text = yaml_content
            
            return {
                "file_name": file_path.name,
                "file_type": "yaml",
                "text": text,
                "word_count": len(text.split()),
                "metadata": {}
            }
            
        except Exception as e:
            raise RuntimeError(f"Failed to parse YAML file: {str(e)}")
    
    @staticmethod
    def parse_xml(file_path: str) -> Dict[str, Any]:
        """Extract content from XML files."""
        if not HAS_HTML:  # BeautifulSoup also handles XML
            raise RuntimeError("XML support not installed. Run: pip install beautifulsoup4 lxml")
        
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                xml_content = f.read()
            
            soup = BeautifulSoup(xml_content, 'lxml-xml')
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            text = '\n'.join(line for line in lines if line)
            
            return {
                "file_name": file_path.name,
                "file_type": "xml",
                "text": text,
                "word_count": len(text.split()),
                "metadata": {}
            }
            
        except Exception as e:
            raise RuntimeError(f"Failed to parse XML file: {str(e)}")
    
    # ==================== MAIN PARSER ====================
    
    @staticmethod
    def parse_file(file_path: str) -> Dict[str, Any]:
        """
        Parse a file based on its extension.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Dictionary containing extracted text and metadata
        """
        file_path = Path(file_path)
        suffix = file_path.suffix.lower()
        
        # Route to appropriate parser
        if suffix == '.docx':
            return DocumentParser.parse_docx(str(file_path))
        elif suffix == '.pdf':
            return DocumentParser.parse_pdf(str(file_path))
        elif suffix == '.rtf':
            return DocumentParser.parse_rtf(str(file_path))
        elif suffix in {'.txt', '.md', '.rst'}:
            return DocumentParser.parse_txt(str(file_path))
        elif suffix == '.pptx':
            return DocumentParser.parse_pptx(str(file_path))
        elif suffix == '.csv':
            return DocumentParser.parse_csv(str(file_path))
        elif suffix in {'.xlsx', '.xls'}:
            return DocumentParser.parse_excel(str(file_path))
        elif suffix in DocumentParser.CODE_FORMATS:
            return DocumentParser.parse_code(str(file_path))
        elif suffix in {'.html', '.htm'}:
            return DocumentParser.parse_html(str(file_path))
        elif suffix == '.json':
            return DocumentParser.parse_json(str(file_path))
        elif suffix in {'.yaml', '.yml'}:
            return DocumentParser.parse_yaml(str(file_path))
        elif suffix == '.xml':
            return DocumentParser.parse_xml(str(file_path))
        else:
            supported = ', '.join(DocumentParser.get_supported_formats())
            raise ValueError(
                f"Unsupported file format: {suffix}\n"
                f"Supported formats: {supported}"
            )
